"""AI Scientist pitch deck — 4 slides, 2-minute format. Rebuilt for platform release."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT     = RGBColor(0x00, 0xC2, 0xFF)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xD6, 0xE0)
GREEN      = RGBColor(0x00, 0xE5, 0x96)
YELLOW     = RGBColor(0xFF, 0xD6, 0x00)
RED_SOFT   = RGBColor(0xFF, 0x6B, 0x6B)
MID_BLUE   = RGBColor(0x12, 0x28, 0x3E)
EMERALD    = RGBColor(0x05, 0xC4, 0x8A)
PURPLE     = RGBColor(0xA8, 0x6E, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── helpers ────────────────────────────────────────────────────────────────────

def add_slide():
    sl = prs.slides.add_slide(BLANK)
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BG
    return sl

def rect(slide, l, t, w, h, colour):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = colour
    s.line.fill.background()
    return s

def txt(slide, l, t, w, h, text, size=16, bold=False, italic=False,
        colour=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = colour
    return tb

def bullets(slide, l, t, w, h, lines, size=14, colour=WHITE, leading=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        text, opts = (line, {}) if not isinstance(line, tuple) else line
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", PP_ALIGN.LEFT)
        if leading:
            p.line_spacing = Pt(leading)
        r = p.add_run(); r.text = text
        r.font.size   = Pt(opts.get("size", size))
        r.font.bold   = opts.get("bold", False)
        r.font.italic = opts.get("italic", False)
        r.font.color.rgb = opts.get("colour", colour)

def hyperlink_txt(slide, l, t, w, h, text, url, size=14, colour=ACCENT,
                  align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = False
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = colour; r.font.underline = True
    rPr = r._r.get_or_add_rPr()
    rId = slide.part.relate_to(
        url,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
        is_external=True)
    hl = etree.SubElement(rPr, qn("a:hlinkClick"))
    hl.set(qn("r:id"), rId)

def header(slide, title):
    rect(slide, 0, 0, 13.33, 0.10, ACCENT)
    rect(slide, 0, 7.40, 13.33, 0.10, ACCENT)
    txt(slide, 0.5, 0.18, 12, 0.60, title, size=32, bold=True, colour=WHITE)
    rect(slide, 0.5, 0.82, 1.2, 0.07, ACCENT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 0.10, ACCENT)
rect(sl, 0, 7.40, 13.33, 0.10, ACCENT)

rect(sl, 0, 0.10, 13.33, 0.42, RGBColor(0x05, 0x12, 0x1E))
txt(sl, 0.5, 0.16, 12.3, 0.32,
    "Hack-Nation 5th Global AI Hackathon  ·  HUB-LINZ",
    size=13, colour=LIGHT_GREY, align=PP_ALIGN.CENTER)

txt(sl, 0.5, 0.75, 12.3, 1.10,
    "Skills4People",
    size=62, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

rect(sl, 3.2, 1.95, 6.93, 0.55, RGBColor(0x00, 0x40, 0x60))
txt(sl, 3.2, 2.00, 6.93, 0.44,
    "The AI Scientist  —  from one-shot tool to experiment platform that learns",
    size=13, italic=True, colour=ACCENT, align=PP_ALIGN.CENTER)

rect(sl, 4.0, 2.68, 5.33, 0.06, ACCENT)

txt(sl, 0.5, 2.88, 12.3, 0.80,
    "GF.sh",
    size=44, bold=True, colour=ACCENT, align=PP_ALIGN.CENTER)

txt(sl, 0.5, 3.78, 12.3, 0.45,
    "Samuel Hajek  ·  Georg Niess  ·  Felix Scope  ·  Johannes Wagner",
    size=17, colour=LIGHT_GREY, align=PP_ALIGN.CENTER)

rect(sl, 3.5, 4.45, 6.33, 0.05, RGBColor(0x22, 0x44, 0x55))
txt(sl, 0.5, 4.60, 12.3, 0.38,
    "OpenAI GPT-4o  ·  Google Gemini  ·  Claude  ·  Supabase pgvector  ·  Tavily  ·  Next.js 14",
    size=12, colour=RGBColor(0x66, 0x99, 0xBB), align=PP_ALIGN.CENTER)

hyperlink_txt(sl, 0.5, 4.95, 12.3, 0.38,
    "🔗  github.com/fescofesco/Skill4People",
    "https://github.com/fescofesco/Skill4People",
    size=14, colour=ACCENT, align=PP_ALIGN.CENTER)

rect(sl, 0.5, 5.45, 12.3, 1.50, RGBColor(0x08, 0x1E, 0x30))
txt(sl, 0.7, 5.60, 11.9, 1.20,
    "Science moves at the speed of operations, not ideas.\n"
    "Every plan saved, every correction made, every document uploaded —\n"
    "the next plan is better.",
    size=16, italic=True, colour=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM + SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 0.10, ACCENT)
rect(sl, 0, 7.40, 13.33, 0.10, ACCENT)

txt(sl, 0.5, 0.22, 12.3, 0.85,
    "THE AI SCIENTIST",
    size=46, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
txt(sl, 0.5, 1.10, 12.3, 0.50,
    "Hypothesis to runnable experiment plan — in minutes, not weeks",
    size=20, italic=True, colour=ACCENT, align=PP_ALIGN.CENTER)
rect(sl, 4.5, 1.70, 4.33, 0.06, ACCENT)

# Problem
rect(sl, 0.4, 1.90, 5.9, 4.90, MID_BLUE)
rect(sl, 0.4, 1.90, 5.9, 0.08, RED_SOFT)
txt(sl, 0.55, 2.03, 5.6, 0.50, "THE PROBLEM", size=18, bold=True, colour=RED_SOFT)
bullets(sl, 0.55, 2.65, 5.6, 3.8,
    ["Turning a hypothesis into a lab-ready experiment\ntakes weeks of manual work",
     "",
     "• Designing protocols from scattered literature",
     "• Sourcing reagents with correct catalog numbers",
     "• Building realistic budgets + phased timelines",
     "• Assessing safety, ethics & regulatory requirements",
     "• Applying your lab's own SOPs and prior corrections",
     "",
     ("The bottleneck isn't ideas — it's operations.", {"italic": True, "colour": YELLOW})],
    size=13, colour=LIGHT_GREY, leading=19)

txt(sl, 6.35, 4.10, 0.6, 0.60, "→", size=36, bold=True, colour=ACCENT,
    align=PP_ALIGN.CENTER)

# Solution
rect(sl, 7.05, 1.90, 5.85, 4.90, RGBColor(0x08, 0x1E, 0x12))
rect(sl, 7.05, 1.90, 5.85, 0.08, GREEN)
txt(sl, 7.20, 2.03, 5.6, 0.50, "OUR SOLUTION", size=18, bold=True, colour=GREEN)
bullets(sl, 7.20, 2.65, 5.6, 3.9,
    [("An iterative experiment platform:", {"bold": True, "colour": WHITE}),
     "",
     ("① Input  — hypothesis + category + optional\n"
      "   reference documents (PDF / TXT)", {"colour": ACCENT}),
     "",
     ("② Literature QC  — 6-source parallel search\n"
      "   → novelty signal in seconds", {"colour": ACCENT}),
     "",
     ("③ Plan + Critique  — full operational plan,\n"
      "   then AI audits its own weaknesses", {"colour": ACCENT}),
     "",
     ("④ Save → Feedback → Regenerate  — corrections\n"
      "   stored by scope, applied to every future plan", {"colour": ACCENT}),
     "",
     ("The platform gets smarter with every use.", {"italic": True, "bold": True, "colour": GREEN})],
    size=12, colour=LIGHT_GREY, leading=17)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHAT WE BUILT
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "WHAT WE BUILT")

features = [
    (RED_SOFT, "🔍  PLAN CRITIQUE — AI EVALUATES ITS OWN WEAKNESSES",
     "After every plan, a second AI pass scans for 6 weakness categories:\n"
     "missing controls · weak statistics · insufficient sample size\n"
     "validation gaps · safety oversights · feasibility issues\n"
     "Rated: weak (critical) · needs_work · solid"),

    (EMERALD,  "🔁  3-BUCKET FEEDBACK PLATFORM  (Stretch Goal ✔)",
     "Corrections scoped to: Organisation · Category · Experiment\n"
     "AI auto-classifies each correction into the right bucket.\n"
     "Stored as Gemini 768-dim vectors in Supabase pgvector (HNSW).\n"
     "Injected as labelled blocks into every new plan. ~1,600 rules pre-seeded."),

    (YELLOW,   "📂  DOCUMENT UPLOAD — YOUR SOPs BECOME PART OF EVERY PLAN",
     "Upload PDF, TXT, or MD at organisation-scope or experiment-scope.\n"
     "Text extracted via pdf-parse (60 KB cap), threaded into plan prompt.\n"
     "Org documents inject into all plans. Experiment docs inject on branch.\n"
     "Your lab's protocols, not just the literature."),

    (PURPLE,   "🏛️  EXPERIMENT LIBRARY + ORGANISATION LAYER",
     "Named experiments saved and re-opened from a persistent library.\n"
     "Save without re-running · Regenerate with latest feedback applied.\n"
     "Category system (cell-biology, diagnostics, …) scopes feedback rules.\n"
     "Multi-tenant org layer: settings drawer, per-org categories & documents."),
]

for i, (colour, title, body) in enumerate(features):
    row, col = divmod(i, 2)
    cx = 0.4  + col * 6.45
    cy = 1.15 + row * 2.90
    rect(sl, cx, cy, 6.1, 2.72, MID_BLUE)
    rect(sl, cx, cy, 6.1, 0.08, colour)
    txt(sl, cx+0.15, cy+0.14, 5.8, 0.48, title, size=13, bold=True, colour=colour)
    txt(sl, cx+0.15, cy+0.66, 5.8, 1.92, body, size=11, colour=LIGHT_GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — BEYOND THE BRIEF
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "WHAT WE DID BETTER THAN THE REQUIREMENTS")

rect(sl, 0.4,  1.05, 5.7, 0.40, RGBColor(0x00, 0x30, 0x55))
rect(sl, 6.3,  1.05, 6.6, 0.40, RGBColor(0x00, 0x40, 0x20))
txt(sl, 0.55,  1.10, 5.5, 0.28, "BRIEF ASKED FOR", size=13, bold=True, colour=LIGHT_GREY)
txt(sl, 6.45,  1.10, 6.3, 0.28, "WE DELIVERED", size=13, bold=True, colour=GREEN)

rows = [
    ("Novelty signal + 1–3 references",
     "6-source parallel search · domain filtering · Jaccard re-ranking · de-duplication"),

    ("Step-by-step protocol",
     "Protocol + AI CRITIC: 6 weakness categories audited after every generation"),

    ("Materials with catalog numbers",
     "Tavily supplier search + regex extraction of prices, pack sizes, catalog numbers"),

    ("Phased timeline",
     "Timeline with decision gates · schedule risks · dependencies per phase"),

    ("[Stretch] Scientist feedback loop",
     "✔  3-bucket RAG: org / category / experiment scope · Gemini embeddings · Supabase pgvector"),

    ("(Not in brief)",
     "Document upload: PDF/TXT/MD · extracted & injected · org-scope + experiment-scope"),

    ("(Not in brief)",
     "Persistent experiment library · named plans · Save + Regenerate workflow"),

    ("(Not in brief)",
     "Category system + multi-tenant org layer · settings drawer · per-org rules"),

    ("(Not in brief)",
     "Multi-model: OpenAI GPT-4o · Google Gemini · Claude — redundancy & best-fit routing"),
]

for i, (left, right) in enumerate(rows):
    y = 1.55 + i * 0.655
    bg = RGBColor(0x10, 0x22, 0x33) if i % 2 == 0 else RGBColor(0x0D, 0x1B, 0x2A)
    rect(sl, 0.4, y, 12.5, 0.63, bg)
    is_beyond = left.startswith("(Not")
    is_stretch = "Stretch" in left
    lcol = RGBColor(0x77, 0x88, 0x99) if is_beyond else LIGHT_GREY
    rcol = EMERALD if is_stretch else (ACCENT if is_beyond else GREEN)
    txt(sl, 0.55, y+0.11, 5.6, 0.44, left,  size=11, colour=lcol, italic=is_beyond)
    txt(sl, 6.30, y+0.11, 6.7, 0.44, right, size=11, colour=rcol)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — FEEDBACK MODEL + AI VALIDATION + CONTACT
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "FEEDBACK MODEL & AI VALIDATION")

# ── LEFT: Feedback model ──────────────────────────────────────────────────────
rect(sl, 0.4, 1.05, 7.0, 4.90, MID_BLUE)
rect(sl, 0.4, 1.05, 7.0, 0.08, EMERALD)
txt(sl, 0.55, 1.18, 6.7, 0.42, "HOW THE FEEDBACK MODEL WORKS", size=14, bold=True, colour=EMERALD)

# Bucket diagram: three stacked boxes with arrows
bucket_data = [
    (EMERALD,  "ORGANISATION",  "Applies to all plans in the org\ne.g. 'Always verify matrix effects for biosensors'"),
    (ACCENT,   "CATEGORY",      "Applies to one experiment type\ne.g. 'For Cell Biology: include mycoplasma testing'"),
    (YELLOW,   "EXPERIMENT",    "Applies when branching from a specific plan\ne.g. 'n=4 sufficient for this HeLa line'"),
]
for i, (col, label, desc) in enumerate(bucket_data):
    by = 1.72 + i * 1.08
    rect(sl, 0.55, by, 6.7, 0.95, RGBColor(0x0A, 0x1E, 0x30))
    rect(sl, 0.55, by, 0.18, 0.95, col)
    txt(sl, 0.85, by+0.08, 2.0, 0.30, label, size=12, bold=True, colour=col)
    txt(sl, 0.85, by+0.40, 6.1, 0.45, desc,  size=11, colour=LIGHT_GREY)

# Storage / retrieval tech line
rect(sl, 0.4, 4.98, 7.0, 0.48, RGBColor(0x06, 0x14, 0x20))
bullets(sl, 0.55, 5.03, 6.7, 0.38,
    [("Gemini 768-dim embeddings · Supabase pgvector HNSW · cosine similarity · "
      "match_experiments() RPC · ~1,600 pre-seeded rules · AI auto-classifies bucket",
      {"colour": RGBColor(0x88, 0xBB, 0xDD), "size": 10})],
    size=10)

# ── RIGHT: AI Validation ──────────────────────────────────────────────────────
rect(sl, 7.6, 1.05, 5.3, 4.90, RGBColor(0x0F, 0x22, 0x35))
rect(sl, 7.6, 1.05, 5.3, 0.08, RED_SOFT)
txt(sl, 7.75, 1.18, 5.0, 0.42, "HOW WE VALIDATE THE AI", size=14, bold=True, colour=RED_SOFT)

validation = [
    (RED_SOFT, "Plan Critique (self-audit)",
     "Second AI pass after every generation.\n6 categories: controls · statistics · sample size\nvalidation · safety · feasibility\nRated: weak · needs_work · solid"),
    (YELLOW,   "Multi-model redundancy",
     "OpenAI GPT-4o · Google Gemini · Claude\nHeuristic fallback at every AI call\nSystem is fully functional with no API key"),
    (ACCENT,   "Composite confidence score",
     "Every plan scored: evidence quality\nsupplier completeness · validation design\nfeedback relevance — visible before ordering"),
    (GREEN,    "Literature grounding",
     "Plans grounded in 6-source literature QC\nNov­elty classified with confidence score\nNo hallucinated catalog numbers — Tavily-verified"),
]
for i, (col, title, body) in enumerate(validation):
    vy = 1.66 + i * 1.05
    rect(sl, 7.75, vy, 4.95, 0.92, RGBColor(0x0A, 0x1A, 0x28))
    rect(sl, 7.75, vy, 4.95, 0.06, col)
    txt(sl, 7.90, vy+0.10, 4.6, 0.28, title, size=11, bold=True, colour=col)
    txt(sl, 7.90, vy+0.40, 4.6, 0.48, body,  size=10, colour=LIGHT_GREY)

# ── CONTACT BAR ───────────────────────────────────────────────────────────────
rect(sl, 0, 6.10, 13.33, 1.30, RGBColor(0x05, 0x12, 0x1E))
rect(sl, 0, 6.10, 13.33, 0.06, ACCENT)

txt(sl, 0.5, 6.18, 12.3, 0.32,
    "CONTACT US FOR MORE QUESTIONS",
    size=13, bold=True, colour=ACCENT, align=PP_ALIGN.CENTER)

contacts = [
    ("Samuel Hajek",    "samuel.hajek@gmail.com"),
    ("Johannes Wagner", "johannes.wagner@alumni.tugraz.at"),
    ("Felix Scope",     "scope@tugraz.at"),
    ("Georg Niess",     "georg.niess@tugraz.at"),
]
for i, (name, email) in enumerate(contacts):
    cx = 0.45 + i * 3.22
    txt(sl, cx, 6.56, 3.1, 0.28, name,  size=11, bold=True,  colour=WHITE)
    txt(sl, cx, 6.84, 3.1, 0.26, email, size=10, italic=True, colour=ACCENT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — CONTACT US
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 0.10, ACCENT)
rect(sl, 0, 7.40, 13.33, 0.10, ACCENT)

txt(sl, 0.5, 0.55, 12.3, 0.80,
    "CONTACT US",
    size=48, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

txt(sl, 0.5, 1.45, 12.3, 0.45,
    "Questions, collaborations, or feedback — we'd love to hear from you.",
    size=16, italic=True, colour=LIGHT_GREY, align=PP_ALIGN.CENTER)

rect(sl, 4.0, 2.05, 5.33, 0.06, ACCENT)

contacts = [
    ("Samuel Hajek",    "Scientist",              "samuel.hajek@gmail.com"),
    ("Johannes Wagner", "Mechanical Engineer",    "johannes.wagner@alumni.tugraz.at"),
    ("Felix Scope",     "Computer Scientist",     "scope@tugraz.at"),
    ("Georg Niess",     "Medicine MSc",           "georg.niess@tugraz.at"),
]

for i, (name, role, email) in enumerate(contacts):
    cx = 0.55 + (i % 2) * 6.4
    cy = 2.35 + (i // 2) * 2.05
    rect(sl, cx, cy, 6.1, 1.75, MID_BLUE)
    rect(sl, cx, cy, 6.1, 0.08, ACCENT)
    txt(sl, cx+0.22, cy+0.18, 5.7, 0.45, name,  size=20, bold=True,  colour=WHITE)
    txt(sl, cx+0.22, cy+0.65, 5.7, 0.30, role,  size=13, italic=True, colour=ACCENT)
    txt(sl, cx+0.22, cy+1.02, 5.7, 0.32, email, size=13, colour=LIGHT_GREY)

hyperlink_txt(sl, 0.5, 6.60, 12.3, 0.38,
    "🔗  github.com/fescofesco/Skill4People",
    "https://github.com/fescofesco/Skill4People",
    size=14, colour=ACCENT, align=PP_ALIGN.CENTER)


# ── save ──────────────────────────────────────────────────────────────────────
out = r"C:\Users\Admin\Documents\Repos\Hacknation_april_2026\presentation\GF_sh_AI_scientist_Deck.pptx"
prs.save(out)
print(f"Saved: {out}")
